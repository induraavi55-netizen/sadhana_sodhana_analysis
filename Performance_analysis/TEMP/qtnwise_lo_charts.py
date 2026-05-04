import pandas as pd
import matplotlib.pyplot as plt
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

DATA_DIR = Path("data")

prs = Presentation()

excel_files = [
    f for f in DATA_DIR.glob("*.xlsx")
    if f.name.lower().startswith("grade") and not f.name.startswith("~$")
]

for excel_file in excel_files:

    xls = pd.ExcelFile(excel_file)

    loq_sheets = [s for s in xls.sheet_names if s.endswith("_lo_question")]

    for loq_sheet in loq_sheets:

        base = loq_sheet.replace("_lo_question", "")
        qlvl_sheet = f"{base}_qlvl"

        if qlvl_sheet not in xls.sheet_names:
            continue

        df_loq = pd.read_excel(excel_file, sheet_name=loq_sheet)
        df_diff = pd.read_excel(excel_file, sheet_name=qlvl_sheet)

        # split difficulty questions
        df_diff["Questions"] = df_diff["Questions"].str.split(",")
        df_diff = df_diff.explode("Questions")
        df_diff = df_diff.rename(columns={"Questions": "Question"})
        df_diff["Question"] = df_diff["Question"].str.strip()

        # merge
        df = df_loq.merge(df_diff[["Question","Difficulty"]], on="Question", how="left")

        color_map = {
            "Low (Content)": "#9ecae1",
            "Medium (Concept)": "#f2b600",
            "High (Challenge)": "#ff1f1f"
        }

        df["color"] = df["Difficulty"].map(color_map)

        df["label"] = df["LO"] + " (" + df["Question"] + ")"

        df = df.sort_values("Performance (%)")

        fig, ax = plt.subplots(figsize=(12,7))

        bars = ax.barh(
            df["label"],
            df["Performance (%)"],
            color=df["color"],
            edgecolor="black",
            height=0.55
        )

        for bar, val in zip(bars, df["Performance (%)"]):
            width = bar.get_width()
            ax.text(
                width + 1,
                bar.get_y() + bar.get_height()/2,
                f"{int(val)}",
                va="center",
                fontsize=10,
                fontweight="bold"
            )

        ax.set_xlabel("Performance %")
        ax.set_ylabel("Learning outcomes")

        title = f"{excel_file.stem} - {base}"
        ax.set_title(title)

        ax.grid(axis="x", linestyle="-", alpha=0.3)

        plt.tight_layout()

        chart_path = DATA_DIR / f"{excel_file.stem}_{base}_chart.png"
        plt.savefig(chart_path, dpi=300, bbox_inches="tight")
        plt.close()

        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.add_picture(
            str(chart_path),
            Inches(0.5),
            Inches(0.8),
            height=Inches(6)
        )

ppt_path = DATA_DIR / "LO_Performance_Report.pptx"
prs.save(ppt_path)

print("All charts exported to PPT.")